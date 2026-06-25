"""文档解析服务 - 支持多种文档格式的解析"""
import os
import json
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


class DocumentParser:
    """文档解析器 - 支持PDF、DOCX、DOC、WPS、XLSX、XLS、PPTX、PPT、TXT、MD、HTML、JSON、图片、代码等格式"""

    # 文件类型映射
    TYPE_MAP = {
        '.pdf': 'pdf',
        '.docx': 'docx', '.doc': 'doc', '.wps': 'wps',
        '.xlsx': 'xlsx', '.xls': 'xls',
        '.pptx': 'pptx', '.ppt': 'ppt',
        '.txt': 'txt', '.md': 'markdown', '.html': 'html', '.htm': 'html',
        '.json': 'json',
        '.jpg': 'image', '.jpeg': 'image', '.png': 'image',
        '.c': 'c_code', '.h': 'c_code',
        '.cpp': 'cpp_code', '.hpp': 'cpp_code', '.cc': 'cpp_code',
        '.py': 'python_code',
        '.java': 'java_code',
        '.cs': 'dotnet_code',
        '.go': 'go_code', '.rs': 'rust_code',
        '.js': 'js_code', '.ts': 'ts_code',
        '.bin': 'binary', '.dat': 'binary', '.hex': 'binary'
    }

    CODE_TYPES = {'c_code', 'cpp_code', 'python_code', 'java_code', 'dotnet_code',
                  'go_code', 'rust_code', 'js_code', 'ts_code'}

    @classmethod
    def get_file_type(cls, filename):
        ext = Path(filename).suffix.lower()
        return cls.TYPE_MAP.get(ext, 'unknown')

    @classmethod
    def is_supported(cls, filename):
        return cls.get_file_type(filename) != 'unknown'

    def parse(self, file_path, file_type=None):
        """解析文档，返回文本内容和结构化数据"""
        if file_type is None:
            file_type = self.get_file_type(file_path)

        try:
            parser_method = getattr(self, f'_parse_{file_type}', None)
            if parser_method:
                return parser_method(file_path)
            else:
                return self._parse_unknown(file_path)
        except Exception as e:
            logger.error(f"解析文件失败: {file_path}, 类型: {file_type}, 错误: {str(e)}")
            return {'text': '', 'structure': {}, 'error': str(e)}

    def _parse_pdf(self, file_path):
        """解析PDF文档"""
        text_parts = []
        structure = {'type': 'pdf', 'pages': []}

        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text() or ''
                    text_parts.append(page_text)
                    structure['pages'].append({
                        'page_num': i + 1,
                        'text_length': len(page_text),
                        'has_tables': bool(page.extract_tables()),
                        'has_images': bool(page.images)
                    })
        except ImportError:
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(file_path)
                for page in reader.pages:
                    text_parts.append(page.extract_text() or '')
                structure['pages'] = [{'page_num': i + 1} for i in range(len(reader.pages))]
            except ImportError:
                logger.warning("未安装PDF解析库，请安装pdfplumber或PyPDF2")

        return {'text': '\n'.join(text_parts), 'structure': structure}

    def _parse_docx(self, file_path):
        """解析DOCX文档"""
        text_parts = []
        structure = {'type': 'docx', 'paragraphs': [], 'tables': []}

        from docx import Document
        doc = Document(file_path)

        for i, para in enumerate(doc.paragraphs):
            text_parts.append(para.text)
            style_name = para.style.name if para.style else ''
            structure['paragraphs'].append({
                'index': i,
                'style': style_name,
                'text_length': len(para.text),
                'is_heading': 'Heading' in style_name or '标题' in style_name
            })

        for i, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            structure['tables'].append({
                'index': i,
                'rows': len(table.rows),
                'cols': len(table.columns),
                'data': table_data
            })
            text_parts.append('\n'.join(['\t'.join(row) for row in table_data]))

        return {'text': '\n'.join(text_parts), 'structure': structure}

    def _parse_doc(self, file_path):
        """解析DOC文档（旧格式）"""
        try:
            import docx2txt
            text = docx2txt.process(file_path)
            return {'text': text, 'structure': {'type': 'doc'}}
        except ImportError:
            return self._parse_docx(file_path)

    def _parse_wps(self, file_path):
        """解析WPS文档"""
        return self._parse_docx(file_path)

    def _parse_xlsx(self, file_path):
        """解析Excel文档"""
        text_parts = []
        structure = {'type': 'xlsx', 'sheets': []}

        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_data = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else '' for cell in row]
                sheet_data.append(row_data)
                text_parts.append('\t'.join(row_data))

            structure['sheets'].append({
                'name': sheet_name,
                'rows': ws.max_row,
                'cols': ws.max_column,
                'data': sheet_data[:100]  # 限制结构化数据大小
            })

        wb.close()
        return {'text': '\n'.join(text_parts), 'structure': structure}

    def _parse_xls(self, file_path):
        """解析XLS文档"""
        return self._parse_xlsx(file_path)

    def _parse_pptx(self, file_path):
        """解析PPTX文档"""
        text_parts = []
        structure = {'type': 'pptx', 'slides': []}

        from pptx import Presentation
        prs = Presentation(file_path)

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    slide_text.append(shape.text)
                    text_parts.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_data = [cell.text for cell in row.cells]
                        text_parts.append('\t'.join(row_data))

            structure['slides'].append({
                'slide_num': i + 1,
                'text_length': len('\n'.join(slide_text)),
                'shape_count': len(slide.shapes)
            })

        return {'text': '\n'.join(text_parts), 'structure': structure}

    def _parse_ppt(self, file_path):
        """解析PPT文档"""
        return self._parse_pptx(file_path)

    def _parse_txt(self, file_path):
        """解析TXT文档"""
        import chardet
        with open(file_path, 'rb') as f:
            raw = f.read()
            result = chardet.detect(raw)
            encoding = result.get('encoding', 'utf-8') or 'utf-8'

        try:
            text = raw.decode(encoding)
        except (UnicodeDecodeError, TypeError):
            text = raw.decode('utf-8', errors='ignore')

        return {'text': text, 'structure': {'type': 'txt', 'length': len(text)}}

    def _parse_markdown(self, file_path):
        """解析Markdown文档"""
        text = self._parse_txt(file_path)['text']
        structure = {'type': 'markdown', 'headings': []}

        import re
        for match in re.finditer(r'^(#{1,6})\s+(.+)$', text, re.MULTILINE):
            level = len(match.group(1))
            title = match.group(2)
            structure['headings'].append({'level': level, 'title': title})

        return {'text': text, 'structure': structure}

    def _parse_html(self, file_path):
        """解析HTML文档"""
        from bs4 import BeautifulSoup

        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()

        soup = BeautifulSoup(html_content, 'lxml')

        # 移除脚本和样式
        for tag in soup(['script', 'style']):
            tag.decompose()

        text = soup.get_text(separator='\n', strip=True)
        structure = {
            'type': 'html',
            'title': soup.title.string if soup.title else '',
            'headings': [h.get_text(strip=True) for h in soup.find_all(['h1', 'h2', 'h3', 'h4'])],
            'links': len(soup.find_all('a')),
            'images': len(soup.find_all('img'))
        }

        return {'text': text, 'structure': structure}

    def _parse_json(self, file_path):
        """解析JSON文档"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        text = json.dumps(data, ensure_ascii=False, indent=2)
        structure = {
            'type': 'json',
            'keys': list(data.keys()) if isinstance(data, dict) else None,
            'length': len(data) if isinstance(data, (list, dict)) else None
        }

        return {'text': text, 'structure': structure}

    def _parse_image(self, file_path):
        """解析图片文档 - OCR提取文字"""
        text = ''
        structure = {'type': 'image', 'path': file_path}

        try:
            from PIL import Image
            img = Image.open(file_path)
            structure['size'] = img.size
            structure['mode'] = img.mode

            # 尝试OCR
            try:
                import pytesseract
                text = pytesseract.image_to_string(img, lang='chi_sim+eng')
            except ImportError:
                logger.info("未安装pytesseract，跳过OCR")
        except ImportError:
            logger.warning("未安装Pillow库")

        return {'text': text, 'structure': structure}

    def _parse_code(self, file_path, language=''):
        """解析代码文件"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            code = f.read()

        structure = {
            'type': 'code',
            'language': language,
            'length': len(code),
            'line_count': code.count('\n') + 1
        }

        # 提取代码结构
        import re
        if language in ('python',):
            # Python: 提取类和函数定义
            classes = re.findall(r'^class\s+(\w+)', code, re.MULTILINE)
            functions = re.findall(r'^def\s+(\w+)', code, re.MULTILINE)
            structure['classes'] = classes
            structure['functions'] = functions
        elif language in ('java',):
            classes = re.findall(r'(?:public|private|protected)?\s*class\s+(\w+)', code)
            methods = re.findall(r'(?:public|private|protected)?\s+\w+\s+(\w+)\s*\(', code)
            structure['classes'] = classes
            structure['methods'] = methods
        elif language in ('c_code', 'cpp_code'):
            functions = re.findall(r'\b\w+\s+(\w+)\s*\([^)]*\)\s*\{', code)
            structure['functions'] = functions
        elif language in ('csharp', 'dotnet_code'):
            classes = re.findall(r'(?:public|private|internal)?\s*class\s+(\w+)', code)
            methods = re.findall(r'(?:public|private|protected)?\s+\w+\s+(\w+)\s*\(', code)
            structure['classes'] = classes
            structure['methods'] = methods

        return {'text': code, 'structure': structure}

    def _parse_c_code(self, file_path):
        return self._parse_code(file_path, 'c')

    def _parse_cpp_code(self, file_path):
        return self._parse_code(file_path, 'cpp')

    def _parse_python_code(self, file_path):
        return self._parse_code(file_path, 'python')

    def _parse_java_code(self, file_path):
        return self._parse_code(file_path, 'java')

    def _parse_dotnet_code(self, file_path):
        return self._parse_code(file_path, 'csharp')

    def _parse_go_code(self, file_path):
        return self._parse_code(file_path, 'go')

    def _parse_rust_code(self, file_path):
        return self._parse_code(file_path, 'rust')

    def _parse_js_code(self, file_path):
        return self._parse_code(file_path, 'javascript')

    def _parse_ts_code(self, file_path):
        return self._parse_code(file_path, 'typescript')

    def _parse_binary(self, file_path):
        """解析二进制文件 - 读取基本信息"""
        file_size = os.path.getsize(file_path)
        return {
            'text': f'[二进制文件: {Path(file_path).name}, 大小: {file_size}字节]',
            'structure': {'type': 'binary', 'size': file_size}
        }

    def _parse_unknown(self, file_path):
        """尝试以文本方式解析未知格式"""
        try:
            return self._parse_txt(file_path)
        except Exception:
            return self._parse_binary(file_path)
